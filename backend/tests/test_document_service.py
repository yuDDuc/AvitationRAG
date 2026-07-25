import pytest
import os
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from app.services.document_service import DocumentService

@pytest.fixture
def document_service():
    """Fixture khởi tạo DocumentService cho các test case."""
    return DocumentService()

@pytest.fixture
def sample_pdf(tmp_path):
    """Tạo file PDF mẫu để test."""
    file_path = os.path.join(tmp_path, "test.pdf")
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 50), "Hello PDF World")
    doc.save(file_path)
    doc.close()
    return file_path

@pytest.fixture
def sample_docx(tmp_path):
    """Tạo file DOCX mẫu để test."""
    file_path = os.path.join(tmp_path, "test.docx")
    doc = Document()
    doc.add_paragraph("Hello Docx World")
    doc.save(file_path)
    return file_path

@pytest.fixture
def sample_pptx(tmp_path):
    """Tạo file PPTX mẫu để test."""
    file_path = os.path.join(tmp_path, "test.pptx")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Hello Pptx World"
    prs.save(file_path)
    return file_path

def test_extract_text_pdf(document_service, sample_pdf):
    """Kiểm tra trích xuất text từ file PDF."""
    text = document_service.extract_text(sample_pdf)
    assert "Hello PDF World" in text

def test_extract_text_docx(document_service, sample_docx):
    """Kiểm tra trích xuất text từ file DOCX."""
    text = document_service.extract_text(sample_docx)
    assert "Hello Docx World" in text

def test_extract_text_pptx(document_service, sample_pptx):
    """Kiểm tra trích xuất text từ file PPTX."""
    text = document_service.extract_text(sample_pptx)
    assert "Hello Pptx World" in text

def test_extract_text_cleaning(document_service, tmp_path):
    """Kiểm tra logic làm sạch text: loại bỏ khoảng trắng dư thừa."""
    file_path = os.path.join(tmp_path, "clean_test.docx")
    doc = Document()
    doc.add_paragraph("  Hello    World  \n\n  ")
    doc.save(file_path)
    
    text = document_service.extract_text(file_path)
    # _clean_text sẽ biến "  Hello    World  \n\n  " thành "Hello World"
    assert text == "Hello World"

def test_chunk_text_basic(document_service):
    """Kiểm tra logic chia nhỏ văn bản cơ bản."""
    text = "Đây là một đoạn văn bản mẫu dùng để kiểm tra tính năng chia nhỏ văn bản của DocumentService. " * 20
    metadata = {"source": "test_source"}
    chunks = document_service.chunk_text(text, metadata, chunk_size=100, chunk_overlap=20)
    
    assert len(chunks) > 0
    for chunk in chunks:
        assert "text" in chunk
        assert "metadata" in chunk
        assert chunk["metadata"]["source"] == "test_source"
        assert len(chunk["text"]) <= 100

def test_chunk_text_empty(document_service):
    """Kiểm tra trường hợp văn bản rỗng."""
    chunks = document_service.chunk_text("", {})
    assert chunks == []

def test_extract_text_file_not_found(document_service):
    """Kiểm tra lỗi khi không tìm thấy file."""
    with pytest.raises(FileNotFoundError):
        document_service.extract_text("non_existent_file.pdf")

def test_extract_text_unsupported_format(document_service, tmp_path):
    """Kiểm tra lỗi khi định dạng file không được hỗ trợ."""
    file_path = os.path.join(tmp_path, "test.txt")
    with open(file_path, "w") as f:
        f.write("test")
    
    with pytest.raises(ValueError, match="không được hỗ trợ"):
        document_service.extract_text(file_path)
