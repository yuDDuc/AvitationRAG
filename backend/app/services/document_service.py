import os
import re
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from typing import Any, List, Dict, Optional
from langchain_text_splitters import RecursiveCharacterTextSplitter

class DocumentService:
    """
    Service xử lý tài liệu: trích xuất văn bản và chia nhỏ (chunking).
    Hỗ trợ các định dạng: PDF, DOCX, PPTX.
    """

    def __init__(self):
        # Khởi tạo bộ chia văn bản với các thiết lập mặc định
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
            separators=["\n\n", "\n", " ", ""]
        )

    def _clean_text(self, text: str) -> str:
        """
        Làm sạch văn bản: loại bỏ khoảng trắng dư thừa và các ký tự không in được.
        
        Args:
            text: Văn bản thô cần làm sạch.
            
        Returns:
            Văn bản đã được làm sạch.
        """
        # Thay thế nhiều khoảng trắng hoặc dòng trống bằng một khoảng trắng duy nhất
        text = re.sub(r'\s+', ' ', text)
        # Loại bỏ các ký tự không in được (control characters)
        text = "".join(char for char in text if char.isprintable())
        return text.strip()

    def extract_text(self, file_path: str) -> List[Dict[str, Any]]: # Thay đổi kiểu trả về
        """
        Trích xuất văn bản từ file PDF, DOCX hoặc PPTX.
        
        Returns:
            Văn bản đã được trích xuất và làm sạch.
            Nếu là PDF, trả về list các dict (text, page). Các định dạng khác trả về dict {text:..., page: N/A}.
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"Không tìm thấy file: {file_path}")

        ext = os.path.splitext(file_path)[1].lower()
        
        extracted_data = [] # Để lưu text và page/slide
        
        try:
            if ext == '.pdf':
                extracted_data = self._extract_from_pdf(file_path)
            elif ext == '.docx':
                text = self._extract_from_docx(file_path)
                extracted_data = [{"text": self._clean_text(text), "page": "N/A"}]
            elif ext == '.pptx':
                text = self._extract_from_pptx(file_path)
                extracted_data = [{"text": self._clean_text(text), "page": "N/A"}]
            else:
                raise ValueError(f"Định dạng file {ext} không được hỗ trợ. Chỉ hỗ trợ PDF, DOCX, PPTX.")
        except Exception as e:
            raise Exception(f"Lỗi khi trích xuất text từ {file_path}: {str(e)}")

        return extracted_data

    def _extract_from_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        """Trích xuất text từ file PDF dùng PyMuPDF, kèm theo số trang."""
        pages_data = []
        with fitz.open(file_path) as doc:
            for page_num, page in enumerate(doc):
                text = page.get_text()
                # Làm sạch text ngay sau khi trích xuất từng trang
                cleaned_text = self._clean_text(text)
                if cleaned_text: # Chỉ thêm trang có nội dung đã làm sạch
                    pages_data.append({"text": cleaned_text, "page": page_num + 1}) # page_num bắt đầu từ 0
        return pages_data

    def _extract_from_docx(self, file_path: str) -> str:
        """Trích xuất text từ file DOCX dùng python-docx."""
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])

    def _extract_from_pptx(self, file_path: str) -> str:
        """Trích xuất text từ file PPTX dùng python-pptx."""
        text = ""
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
        return text

    def chunk_text(self, extracted_data: List[Dict[str, Any]], common_metadata: Dict, chunk_size: Optional[int] = None, chunk_overlap: Optional[int] = None) -> List[Dict]:
        """
        Chia nhỏ văn bản thành các đoạn (chunks) kèm theo metadata.
        
        Args:
            extracted_data: List các dict chứa 'text' và 'page' (hoặc các metadata khác).
            common_metadata: Thông tin bổ sung chung cho toàn bộ tài liệu (ví dụ: tên file, môn học).
        """
        if not extracted_data:
            return []

        all_chunks = []
        splitter = self.text_splitter
        if chunk_size is not None or chunk_overlap is not None:
            splitter = RecursiveCharacterTextSplitter(
                chunk_size=chunk_size or self.text_splitter._chunk_size,
                chunk_overlap=chunk_overlap or self.text_splitter._chunk_overlap,
                length_function=len,
                separators=["\n\n", "\n", " ", ""]
            )
        
        for item in extracted_data:
            text_to_chunk = item["text"]
            page_metadata = {"page": item.get("page", "N/A")}
            
            # Nếu text quá nhỏ thì không cần chia nhỏ thêm
            if len(text_to_chunk) <= (chunk_size or self.text_splitter._chunk_size):
                chunk_metadata = {**common_metadata, **page_metadata}
                all_chunks.append({"text": text_to_chunk, "metadata": chunk_metadata})
                continue

            chunks_from_page = splitter.split_text(text_to_chunk)
            
            for chunk in chunks_from_page:
                chunk_metadata = {**common_metadata, **page_metadata}
                all_chunks.append({
                    "text": chunk,
                    "metadata": chunk_metadata
                })
        
        return all_chunks
